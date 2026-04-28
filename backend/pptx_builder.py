"""Сборка PPTX по плану презентации.

Все слайды рендерятся через python-pptx с кастомным цветовым оформлением,
без зависимости от шаблонов PowerPoint.
"""
from __future__ import annotations

import io
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .design_presets import merge_slide_palette, style_for_preset
from .slide_planner import PresentationPlan, SlideSpec

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _slide_palette(plan: PresentationPlan, spec: SlideSpec) -> dict[str, str]:
    return merge_slide_palette(plan.palette, getattr(spec, "style", None) or {})


def _preset_pid(plan: PresentationPlan) -> str:
    return getattr(plan, "design_preset", None) or "fresh"


def _preset_style(plan: PresentationPlan) -> dict[str, str | float]:
    return style_for_preset(_preset_pid(plan))


def _decorate_modern_title(slide, pal: dict[str, str], pid: str) -> None:
    """Декор титульного слайда: лёгкий акцент справа сверху."""
    a2 = pal.get("accent2", pal["accent"])
    orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(0.35), Inches(2.9), Inches(2.9))
    orb.fill.solid()
    orb.fill.fore_color.rgb = _hex(a2)
    orb.line.fill.background()
    orb.shadow.inherit = False
    if pid != "midnight":
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(4.2), Inches(0.14))
        _set_solid_fill(stripe, pal["accent"])
        stripe.shadow.inherit = False


def _decorate_content_header(slide, pal: dict[str, str], pid: str) -> None:
    """Шапка контентного слайда: тонкая верхняя линия акцентом, пятно accent2."""
    a2 = pal.get("accent2", pal["accent"])
    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    _set_solid_fill(top_band, a2)
    top_band.shadow.inherit = False
    if pid not in ("midnight",):
        spot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(0.35), Inches(1.9), Inches(1.9))
        spot.fill.solid()
        spot.fill.fore_color.rgb = _hex(a2)
        spot.line.fill.background()


def _accent_bar_width(pid: str):
    return Inches(0.22) if pid in ("ocean", "fresh", "sunrise") else Inches(0.16)


def _hex(value: str) -> RGBColor:
    v = value.lstrip("#")
    if len(v) != 6:
        v = "111827"
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _set_solid_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color)
    shape.line.fill.background()


def _add_background(slide, color: str) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_solid_fill(rect, color)
    rect.shadow.inherit = False


def _add_accent_bar(slide, color: str, x=0, y=0, w=Inches(0.18), h=SLIDE_H) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_solid_fill(bar, color)


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    color: str,
    size: int,
    bold: bool = False,
    align: str = "left",
    font_name: str = "Calibri",
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.name = font_name or "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _hex(color)


def _add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets: list[str],
    *,
    color: str,
    accent: str,
    size: int = 20,
    font_name: str = "Calibri",
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        marker = p.add_run()
        marker.text = "●  "
        marker.font.name = font_name or "Calibri"
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = _hex(accent)

        run = p.add_run()
        run.text = item
        run.font.name = font_name or "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = _hex(color)


def _add_image(slide, image_bytes: bytes, left, top, width, height) -> None:
    try:
        with PILImage.open(io.BytesIO(image_bytes)) as im:
            iw, ih = im.size
    except Exception:
        return

    box_w = int(width)
    box_h = int(height)
    if iw <= 0 or ih <= 0 or box_w <= 0 or box_h <= 0:
        return

    ratio = min(box_w / iw, box_h / ih)
    draw_w = int(iw * ratio)
    draw_h = int(ih * ratio)
    draw_left = int(left) + (box_w - draw_w) // 2
    draw_top = int(top) + (box_h - draw_h) // 2

    bio = io.BytesIO(image_bytes)
    pic = slide.shapes.add_picture(
        bio,
        Emu(draw_left),
        Emu(draw_top),
        width=Emu(draw_w),
        height=Emu(draw_h),
    )
    pic.shadow.inherit = False


def _add_table(
    slide,
    left,
    top,
    width,
    height,
    headers: list[str],
    rows: list[list[str]],
    *,
    palette: dict[str, str],
    font_name: str = "Calibri",
) -> None:
    if not headers:
        return
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 на шапку

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Шапка
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex(palette["primary"])
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(header)
        run.font.name = font_name or "Calibri"
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = _hex(palette["background"])

    # Данные
    body_text = _hex(palette["text"])
    band_color_a = _hex(palette["background"])
    band_color_b = _hex("#F4F6FA")
    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = band_color_a if r % 2 == 1 else band_color_b
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.03)
            tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(row[c]) if c < len(row) else ""
            run.font.name = font_name or "Calibri"
            run.font.size = Pt(13)
            run.font.color.rgb = body_text


def _render_title(prs: Presentation, plan: PresentationPlan, spec: SlideSpec) -> None:
    pid = _preset_pid(plan)
    style = _preset_style(plan)
    title_font = str(style.get("title_font", "Calibri"))
    body_font = str(style.get("body_font", "Calibri"))
    title_scale = float(style.get("title_scale", 1.0))
    body_scale = float(style.get("body_scale", 1.0))
    pal = _slide_palette(plan, spec)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_background(slide, pal["background"])
    _decorate_modern_title(slide, pal, pid)
    _add_accent_bar(slide, pal["accent"], 0, 0, SLIDE_W, Inches(0.22))
    _add_accent_bar(
        slide, pal["accent"], 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22)
    )

    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.4),
        SLIDE_W - Inches(1.8),
        Inches(2.0),
        spec.title or plan.title,
        color=pal["primary"],
        size=max(34, int(54 * title_scale)),
        bold=True,
        align="center",
        font_name=title_font,
    )
    if spec.subtitle or plan.subtitle:
        _add_textbox(
            slide,
            Inches(0.9),
            Inches(4.4),
            SLIDE_W - Inches(1.8),
            Inches(1.2),
            spec.subtitle or plan.subtitle,
            color=pal["muted"],
            size=max(16, int(24 * body_scale)),
            align="center",
            font_name=body_font,
        )


def _render_section(prs: Presentation, plan: PresentationPlan, spec: SlideSpec) -> None:
    style = _preset_style(plan)
    title_font = str(style.get("title_font", "Calibri"))
    body_font = str(style.get("body_font", "Calibri"))
    title_scale = float(style.get("title_scale", 1.0))
    body_scale = float(style.get("body_scale", 1.0))
    pal = _slide_palette(plan, spec)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_background(slide, pal["primary"])
    _add_textbox(
        slide,
        Inches(1.0),
        Inches(3.0),
        SLIDE_W - Inches(2.0),
        Inches(1.6),
        spec.title,
        color=pal["background"],
        size=max(30, int(48 * title_scale)),
        bold=True,
        align="center",
        font_name=title_font,
    )
    if spec.subtitle:
        _add_textbox(
            slide,
            Inches(1.0),
            Inches(4.4),
            SLIDE_W - Inches(2.0),
            Inches(1.0),
            spec.subtitle,
            color=pal["accent"],
            size=max(16, int(24 * body_scale)),
            align="center",
            font_name=body_font,
        )


def _render_content(
    prs: Presentation,
    plan: PresentationPlan,
    spec: SlideSpec,
    image: bytes | None,
) -> None:
    pid = _preset_pid(plan)
    style = _preset_style(plan)
    title_font = str(style.get("title_font", "Calibri"))
    body_font = str(style.get("body_font", "Calibri"))
    title_scale = float(style.get("title_scale", 1.0))
    body_scale = float(style.get("body_scale", 1.0))
    underline_ratio = float(style.get("underline_ratio", 0.09))
    pal = _slide_palette(plan, spec)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_background(slide, pal["background"])
    _decorate_content_header(slide, pal, pid)
    wbar = _accent_bar_width(pid)
    _add_accent_bar(slide, pal["accent"], 0, 0, wbar, SLIDE_H)

    title_left = Inches(0.7)
    title_top = Inches(0.5)
    _add_textbox(
        slide,
        title_left,
        title_top,
        SLIDE_W - Inches(1.4),
        Inches(0.9),
        spec.title,
        color=pal["primary"],
        size=max(24, int(32 * title_scale)),
        bold=True,
        font_name=title_font,
    )
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left,
        Inches(1.45),
        Inches(max(0.9, min(2.1, 15 * underline_ratio))),
        Emu(38100),
    )
    _set_solid_fill(underline, pal["accent"])

    has_image = image is not None
    left_img = has_image and getattr(spec, "image_placement", "right") == "left"
    text_left = Inches(0.7)
    text_top = Inches(1.7)
    text_height = SLIDE_H - text_top - Inches(0.5)
    text_width = SLIDE_W - Inches(1.4)
    img_left = Inches(7.4)
    img_top = Inches(1.65)
    img_w = SLIDE_W - img_left - Inches(0.5)
    img_h = SLIDE_H - img_top - Inches(0.5)
    if has_image:
        if left_img:
            img_left = Inches(0.65)
            img_top = Inches(1.65)
            img_w = Inches(6.35)
            img_h = SLIDE_H - img_top - Inches(0.5)
            text_left = Inches(7.45)
            text_width = SLIDE_W - text_left - Inches(0.55)
        else:
            text_width = Inches(6.4)
            img_left = Inches(7.4)
            img_top = Inches(1.7)
            img_w = SLIDE_W - img_left - Inches(0.5)
            img_h = SLIDE_H - img_top - Inches(0.5)

    if spec.bullets:
        _add_bullets(
            slide,
            text_left,
            text_top,
            text_width,
            text_height,
            spec.bullets,
            color=pal["text"],
            accent=pal["accent"],
            size=max(14, int(20 * body_scale)),
            font_name=body_font,
        )
    elif spec.body:
        _add_textbox(
            slide,
            text_left,
            text_top,
            text_width,
            text_height,
            spec.body,
            color=pal["text"],
            size=max(13, int(18 * body_scale)),
            font_name=body_font,
        )

    if has_image:
        _add_image(slide, image, img_left, img_top, img_w, img_h)


def _render_two_column(
    prs: Presentation,
    plan: PresentationPlan,
    spec: SlideSpec,
    image: bytes | None,
) -> None:
    pid = _preset_pid(plan)
    style = _preset_style(plan)
    title_font = str(style.get("title_font", "Calibri"))
    body_font = str(style.get("body_font", "Calibri"))
    title_scale = float(style.get("title_scale", 1.0))
    body_scale = float(style.get("body_scale", 1.0))
    pal = _slide_palette(plan, spec)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_background(slide, pal["background"])
    _decorate_content_header(slide, pal, pid)
    wbar = _accent_bar_width(pid)
    _add_accent_bar(slide, pal["accent"], 0, 0, wbar, SLIDE_H)

    _add_textbox(
        slide,
        Inches(0.7),
        Inches(0.5),
        SLIDE_W - Inches(1.4),
        Inches(0.9),
        spec.title,
        color=pal["primary"],
        size=max(24, int(32 * title_scale)),
        bold=True,
        font_name=title_font,
    )

    half = (SLIDE_W - Inches(2.1)) / 2
    bullets = spec.bullets or []
    left_b = bullets[: max(1, len(bullets) // 2)] or bullets
    right_b = bullets[len(left_b) :]

    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        half,
        SLIDE_H - Inches(2.2),
        left_b,
        color=pal["text"],
        accent=pal["accent"],
        size=max(14, int(20 * body_scale)),
        font_name=body_font,
    )
    if right_b:
        _add_bullets(
            slide,
            Inches(0.7) + half + Inches(0.7),
            Inches(1.7),
            half,
            SLIDE_H - Inches(2.2),
            right_b,
            color=pal["text"],
            accent=pal["accent"],
            size=max(14, int(20 * body_scale)),
            font_name=body_font,
        )
    elif image is not None:
        _add_image(
            slide,
            image,
            Inches(0.7) + half + Inches(0.7),
            Inches(1.7),
            half,
            SLIDE_H - Inches(2.2),
        )


def _render_table_slide(prs: Presentation, plan: PresentationPlan, spec: SlideSpec) -> None:
    pid = _preset_pid(plan)
    style = _preset_style(plan)
    title_font = str(style.get("title_font", "Calibri"))
    body_font = str(style.get("body_font", "Calibri"))
    title_scale = float(style.get("title_scale", 1.0))
    body_scale = float(style.get("body_scale", 1.0))
    underline_ratio = float(style.get("underline_ratio", 0.09))
    pal = _slide_palette(plan, spec)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_background(slide, pal["background"])
    _decorate_content_header(slide, pal, pid)
    wbar = _accent_bar_width(pid)
    _add_accent_bar(slide, pal["accent"], 0, 0, wbar, SLIDE_H)

    title_left = Inches(0.7)
    _add_textbox(
        slide,
        title_left,
        Inches(0.5),
        SLIDE_W - Inches(1.4),
        Inches(0.9),
        spec.title,
        color=pal["primary"],
        size=max(24, int(32 * title_scale)),
        bold=True,
        font_name=title_font,
    )
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left,
        Inches(1.45),
        Inches(max(0.9, min(2.1, 15 * underline_ratio))),
        Emu(38100),
    )
    _set_solid_fill(underline, pal["accent"])

    if spec.subtitle:
        _add_textbox(
            slide,
            title_left,
            Inches(1.55),
            SLIDE_W - Inches(1.4),
            Inches(0.5),
            spec.subtitle,
            color=pal["muted"],
            size=max(12, int(15 * body_scale)),
            font_name=body_font,
        )

    table_top = Inches(2.2) if spec.subtitle else Inches(1.85)
    table_height = SLIDE_H - table_top - Inches(0.6)
    table_width = SLIDE_W - Inches(1.4)

    if spec.headers and spec.rows:
        _add_table(
            slide,
            title_left,
            table_top,
            table_width,
            table_height,
            spec.headers,
            spec.rows,
            palette=pal,
            font_name=body_font,
        )


def _render_conclusion(
    prs: Presentation, plan: PresentationPlan, spec: SlideSpec
) -> None:
    style = _preset_style(plan)
    title_font = str(style.get("title_font", "Calibri"))
    body_font = str(style.get("body_font", "Calibri"))
    title_scale = float(style.get("title_scale", 1.0))
    body_scale = float(style.get("body_scale", 1.0))
    pal = _slide_palette(plan, spec)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_background(slide, pal["background"])
    _add_accent_bar(slide, pal["accent"], 0, 0, SLIDE_W, Inches(0.22))
    a2 = pal.get("accent2", pal["accent"])
    spot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(0.45), Inches(1.6), Inches(1.6))
    spot.fill.solid()
    spot.fill.fore_color.rgb = _hex(a2)
    spot.line.fill.background()

    _add_textbox(
        slide,
        Inches(0.9),
        Inches(0.9),
        SLIDE_W - Inches(1.8),
        Inches(1.2),
        spec.title or "Выводы",
        color=pal["primary"],
        size=max(28, int(40 * title_scale)),
        bold=True,
        align="center",
        font_name=title_font,
    )

    if spec.bullets:
        _add_bullets(
            slide,
            Inches(1.5),
            Inches(2.4),
            SLIDE_W - Inches(3.0),
            Inches(4.0),
            spec.bullets,
            color=pal["text"],
            accent=pal["accent"],
            size=max(16, int(22 * body_scale)),
            font_name=body_font,
        )
    elif spec.body:
        _add_textbox(
            slide,
            Inches(1.5),
            Inches(2.4),
            SLIDE_W - Inches(3.0),
            Inches(4.0),
            spec.body,
            color=pal["text"],
            size=max(16, int(22 * body_scale)),
            align="center",
            font_name=body_font,
        )

    if spec.subtitle:
        _add_textbox(
            slide,
            Inches(0.9),
            SLIDE_H - Inches(1.0),
            SLIDE_W - Inches(1.8),
            Inches(0.6),
            spec.subtitle,
            color=pal["muted"],
            size=max(13, int(18 * body_scale)),
            align="center",
            font_name=body_font,
        )


def build_pptx(
    plan: PresentationPlan,
    images: dict[int, bytes],
    out_path: Path,
) -> Path:
    """Собрать PPTX. images — словарь {номер_слайда: bytes_картинки}."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, spec in enumerate(plan.slides):
        img = images.get(idx)
        if spec.kind == "title":
            _render_title(prs, plan, spec)
        elif spec.kind == "section":
            _render_section(prs, plan, spec)
        elif spec.kind == "two_column":
            _render_two_column(prs, plan, spec, img)
        elif spec.kind == "conclusion":
            _render_conclusion(prs, plan, spec)
        elif spec.kind == "table":
            _render_table_slide(prs, plan, spec)
        else:
            _render_content(prs, plan, spec, img)

        if spec.notes:
            slide = prs.slides[idx]
            slide.notes_slide.notes_text_frame.text = spec.notes

    prs.save(str(out_path))
    return out_path
