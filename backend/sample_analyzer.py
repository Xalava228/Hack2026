"""Анализ существующей презентации (PPTX или PDF) для извлечения стиля.

Возвращает SampleAnalysis: количество слайдов, оценка плотности текста,
извлечённая цветовая палитра, типы слайдов и краткое содержание.
"""
from __future__ import annotations

import json
import logging
import re
from collections import Counter
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation
from pptx.util import Emu

logger = logging.getLogger(__name__)

SourceFormat = Literal["pptx", "pdf", "docx"]


@dataclass
class SampleSlideInfo:
    title: str = ""
    bullets: list[str] = field(default_factory=list)
    body: str = ""
    word_count: int = 0
    has_image: bool = False
    kind_guess: str = "content"  # title | content | section | conclusion

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass
class SampleAnalysis:
    sample_id: str
    source_format: SourceFormat
    file_name: str
    n_slides: int
    palette: dict[str, str]
    density: str  # minimal | balanced | detailed
    has_images: bool
    slides: list[SampleSlideInfo]
    title_guess: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "sample_id": self.sample_id,
            "source_format": self.source_format,
            "file_name": self.file_name,
            "n_slides": self.n_slides,
            "palette": self.palette,
            "density": self.density,
            "has_images": self.has_images,
            "title_guess": self.title_guess,
            "slides": [s.to_dict() for s in self.slides],
        }

    def short_summary(self) -> dict[str, Any]:
        """Компактная сводка для UI (без полных текстов)."""
        return {
            "sample_id": self.sample_id,
            "source_format": self.source_format,
            "file_name": self.file_name,
            "n_slides": self.n_slides,
            "palette": self.palette,
            "density": self.density,
            "has_images": self.has_images,
            "title_guess": self.title_guess,
            "outline": [
                {
                    "title": s.title[:80],
                    "kind": s.kind_guess,
                    "bullets": len(s.bullets),
                    "has_image": s.has_image,
                }
                for s in self.slides
            ],
        }


_DEFAULT_PALETTE = {
    "primary": "#1F2937",
    "accent": "#6366F1",
    "background": "#FFFFFF",
    "text": "#111827",
    "muted": "#6B7280",
}


def _density_from_words(avg_words: float) -> str:
    if avg_words < 18:
        return "minimal"
    if avg_words < 50:
        return "balanced"
    return "detailed"


def _luminance(rgb: tuple[int, int, int]) -> float:
    r, g, b = (c / 255 for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _hex_to_rgb(s: str) -> tuple[int, int, int]:
    s = s.lstrip("#")
    return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)


def _rgb_to_hex(rgb: tuple[int, int, int]) -> str:
    return "#{:02X}{:02X}{:02X}".format(*rgb)


def _palette_from_counts(counter: Counter[tuple[int, int, int]]) -> dict[str, str]:
    """Из набора цветов собираем согласованную палитру."""
    if not counter:
        return dict(_DEFAULT_PALETTE)

    candidates = [rgb for rgb, _ in counter.most_common(20)]
    candidates.sort(key=_luminance)

    bg = candidates[-1]
    text = candidates[0]
    if _luminance(bg) - _luminance(text) < 0.4:
        bg = (255, 255, 255)
        text = (17, 24, 39)

    middles = [c for c in candidates if c not in (bg, text)]
    accent = next(
        (
            c
            for c in middles
            if 0.18 < _luminance(c) < 0.78
            and (max(c) - min(c)) > 25
        ),
        (99, 102, 241),
    )
    primary = next(
        (c for c in middles if c not in (accent,) and _luminance(c) < 0.45),
        (31, 41, 55),
    )
    muted = next(
        (
            c
            for c in middles
            if c not in (accent, primary) and 0.30 < _luminance(c) < 0.70
        ),
        (107, 114, 128),
    )

    return {
        "primary": _rgb_to_hex(primary),
        "accent": _rgb_to_hex(accent),
        "background": _rgb_to_hex(bg),
        "text": _rgb_to_hex(text),
        "muted": _rgb_to_hex(muted),
    }


def _guess_kind(idx: int, total: int, info: SampleSlideInfo) -> str:
    if idx == 0:
        return "title"
    if idx == total - 1 and re.search(
        r"(вывод|итог|заключени|conclusion|summary|спасибо|thank)",
        (info.title + " " + info.body).lower(),
    ):
        return "conclusion"
    if idx == total - 1:
        return "conclusion"
    if not info.bullets and len(info.body) < 60 and info.word_count <= 8:
        return "section"
    return "content"


# ============================ PPTX ============================
def _extract_pptx_palette(prs: Presentation) -> Counter[tuple[int, int, int]]:
    counter: Counter[tuple[int, int, int]] = Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                color = run.font.color
                                rgb = getattr(color, "rgb", None)
                                if rgb is not None:
                                    counter[(rgb[0], rgb[1], rgb[2])] += 2
                            except Exception:
                                pass
            except Exception:
                pass
            try:
                fill = getattr(shape, "fill", None)
                if fill and getattr(fill, "type", None) == 1:
                    rgb = fill.fore_color.rgb
                    if rgb is not None:
                        counter[(rgb[0], rgb[1], rgb[2])] += 1
            except Exception:
                pass
    return counter


def _extract_pptx_slide(slide) -> SampleSlideInfo:
    title = ""
    bullets: list[str] = []
    body_parts: list[str] = []
    has_image = False

    for shape in slide.shapes:
        try:
            shape_type_name = type(shape).__name__.lower()
            if "picture" in shape_type_name or getattr(shape, "shape_type", None) == 13:
                has_image = True
        except Exception:
            pass

        try:
            if shape.has_text_frame:
                texts = [p.text.strip() for p in shape.text_frame.paragraphs]
                texts = [t for t in texts if t]
                if not texts:
                    continue
                is_title = False
                try:
                    ph = getattr(shape, "placeholder_format", None)
                    if ph is not None and ph.idx in (0, 13):
                        is_title = True
                except Exception:
                    pass
                if not title and is_title:
                    title = texts[0]
                    if len(texts) > 1:
                        body_parts.extend(texts[1:])
                else:
                    if len(texts) > 1 or (
                        texts and len(texts[0]) > 0 and (
                            texts[0].startswith(("•", "-", "·", "●", "*"))
                            or len(texts) >= 2
                        )
                    ):
                        for t in texts:
                            cleaned = t.lstrip("•-·●*— ").strip()
                            if cleaned:
                                if 2 <= len(cleaned.split()) <= 30:
                                    bullets.append(cleaned)
                                else:
                                    body_parts.append(cleaned)
                    else:
                        body_parts.extend(texts)
        except Exception:
            continue

    if not title and body_parts:
        first = body_parts[0]
        if len(first) <= 100:
            title = first
            body_parts = body_parts[1:]

    body = " ".join(body_parts).strip()
    info = SampleSlideInfo(
        title=title,
        bullets=bullets[:10],
        body=body[:600],
        word_count=len((title + " " + body + " " + " ".join(bullets)).split()),
        has_image=has_image,
    )
    return info


def analyze_pptx(path: Path, sample_id: str) -> SampleAnalysis:
    prs = Presentation(str(path))
    slides: list[SampleSlideInfo] = []
    for slide in prs.slides:
        slides.append(_extract_pptx_slide(slide))
    n = len(slides)
    for i, info in enumerate(slides):
        info.kind_guess = _guess_kind(i, n, info)

    palette_counter = _extract_pptx_palette(prs)
    palette = _palette_from_counts(palette_counter)

    avg_words = sum(s.word_count for s in slides) / max(1, n)
    density = _density_from_words(avg_words)
    has_images = any(s.has_image for s in slides)
    title_guess = slides[0].title if slides else ""

    return SampleAnalysis(
        sample_id=sample_id,
        source_format="pptx",
        file_name=path.name,
        n_slides=n,
        palette=palette,
        density=density,
        has_images=has_images,
        slides=slides,
        title_guess=title_guess,
    )


# ============================ PDF ============================
def _split_pdf_page_text(text: str) -> tuple[str, list[str], str]:
    """Эвристика: первая короткая строка = заголовок.

    Дальше — буллеты (короткие/со маркерами) или body (длинный текст)."""
    raw_lines = [ln.rstrip() for ln in text.splitlines()]
    lines = [ln.strip() for ln in raw_lines if ln.strip()]
    if not lines:
        return "", [], ""

    title = ""
    body_lines: list[str] = []
    bullets: list[str] = []

    if len(lines[0]) <= 90:
        title = lines[0]
        rest = lines[1:]
    else:
        rest = lines

    bullet_re = re.compile(r"^[\-\*\u2022\u00b7\u25cf\u25e6\u2013\u2014\u2192]\s+")
    for ln in rest:
        if bullet_re.match(ln):
            cleaned = bullet_re.sub("", ln).strip()
            if cleaned:
                bullets.append(cleaned)
        elif 2 <= len(ln.split()) <= 25 and not ln.endswith(".") and len(rest) > 1:
            bullets.append(ln)
        else:
            body_lines.append(ln)

    body = " ".join(body_lines).strip()
    if not bullets and body:
        sentences = re.split(r"(?<=[.!?])\s+", body)
        if len(sentences) >= 3:
            bullets = [s.strip() for s in sentences[:6] if 3 <= len(s.split()) <= 25]
            if bullets:
                body = ""

    return title, bullets[:10], body[:600]


def analyze_pdf(path: Path, sample_id: str) -> SampleAnalysis:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    slides: list[SampleSlideInfo] = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        title, bullets, body = _split_pdf_page_text(text)
        info = SampleSlideInfo(
            title=title,
            bullets=bullets,
            body=body,
            word_count=len((title + " " + body + " " + " ".join(bullets)).split()),
            has_image=False,
        )
        slides.append(info)

    n = len(slides) or 1
    for i, info in enumerate(slides):
        info.kind_guess = _guess_kind(i, n, info)

    avg_words = sum(s.word_count for s in slides) / max(1, n)
    density = _density_from_words(avg_words)
    title_guess = slides[0].title if slides else ""

    return SampleAnalysis(
        sample_id=sample_id,
        source_format="pdf",
        file_name=path.name,
        n_slides=n,
        palette=dict(_DEFAULT_PALETTE),  # из PDF цвета по тексту извлечь надёжно нельзя
        density=density,
        has_images=False,
        slides=slides,
        title_guess=title_guess,
    )


def _is_heading_style(style_name: str) -> bool:
    n = (style_name or "").strip().lower()
    if not n:
        return False
    if "heading" in n or "заголовок" in n:
        return True
    if n in ("title", "название", "subtitle"):
        return True
    return False


def _paragraphs_from_docx_tables(doc) -> list[tuple[str, str]]:
    """Параграфы и строки таблиц в порядке следования в документе."""
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    rows: list[tuple[str, str]] = []
    for block in doc.element.body:
        if block.tag == qn("w:tbl"):
            table = Table(block, doc)
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    rows.append(("p", line))
        elif block.tag == qn("w:p"):
            para = Paragraph(block, doc)
            t = para.text.strip()
            if not t:
                continue
            st = para.style.name if para.style else ""
            kind = "h" if _is_heading_style(st) else "p"
            rows.append((kind, t))
    return rows


def _blocks_from_docx(doc) -> list[tuple[str, str]]:
    """Параграфы и таблицы в порядке следования в документе."""
    try:
        return _paragraphs_from_docx_tables(doc)
    except Exception:
        logger.exception("docx: fallback to paragraphs only")
        out: list[tuple[str, str]] = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            st = para.style.name if para.style else ""
            kind = "h" if _is_heading_style(st) else "p"
            out.append((kind, t))
        return out


def _docx_section_to_slide(section_title: str, paras: list[str]) -> SampleSlideInfo:
    title = section_title.strip()
    body_paras = list(paras)
    if not title and body_paras:
        first = body_paras[0].strip()
        if len(first) <= 120:
            title = first[:200]
            body_paras = body_paras[1:]
        else:
            title = "Введение"
    if not title:
        title = "Раздел"
    body_join = " ".join(body_paras).strip()
    bullets = [p[:200] for p in body_paras[:12] if len(p.split()) <= 40]
    return SampleSlideInfo(
        title=title[:200],
        bullets=bullets[:10],
        body=body_join[:1200],
        word_count=len((title + " " + body_join).split()),
        has_image=False,
    )


def _extract_docx_paragraphs_zip(path: Path) -> list[str]:
    """Извлекает текст абзацев из DOCX без python-docx (только стандартная библиотека).

    Уступает python-docx в распознавании стилей заголовков и части разметки,
    но позволяет разбирать файл, если пакет не установлен в окружении сервера.
    """
    import zipfile
    import xml.etree.ElementTree as ET

    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    q = lambda tag: f"{{{ns}}}{tag}"

    try:
        with zipfile.ZipFile(path) as z:
            xml_bytes = z.read("word/document.xml")
    except KeyError as e:
        raise ValueError("DOCX без word/document.xml") from e
    except zipfile.BadZipFile as e:
        raise ValueError("Файл не похож на корректный DOCX (ZIP)") from e

    root = ET.fromstring(xml_bytes)
    out: list[str] = []
    for p_el in root.iter(q("p")):
        parts: list[str] = []
        for node in p_el.iter(q("t")):
            if node.text:
                parts.append(node.text)
            if node.tail:
                parts.append(node.tail)
        line = "".join(parts).strip()
        if line:
            out.append(line)
    return out


def _chunk_plain_paragraphs(paragraphs: list[str], target_words: int = 320) -> list[str]:
    """Разбивает сплошной текст на куски по ~target_words слов."""
    chunks: list[str] = []
    current: list[str] = []
    wc = 0
    for line in paragraphs:
        line_w = len(line.split())
        if wc + line_w > target_words and current:
            chunks.append("\n".join(current))
            current = [line]
            wc = line_w
        else:
            current.append(line)
            wc += line_w
    if current:
        chunks.append("\n".join(current))
    return chunks


def analyze_docx(path: Path, sample_id: str) -> SampleAnalysis:
    blocks: list[tuple[str, str]]
    try:
        from docx import Document

        doc = Document(str(path))
        blocks = _blocks_from_docx(doc)
    except ImportError:
        logger.warning(
            "Пакет python-docx не найден; для DOCX используется упрощённое извлечение текста. "
            "Рекомендуется: pip install python-docx"
        )
        plain = _extract_docx_paragraphs_zip(path)
        blocks = [("p", t) for t in plain]

    heading_count = sum(1 for k, _ in blocks if k == "h")
    slides: list[SampleSlideInfo] = []

    if heading_count > 0:
        cur_title = ""
        cur_paras: list[str] = []
        for kind, text in blocks:
            if kind == "h":
                if cur_title or cur_paras:
                    slides.append(_docx_section_to_slide(cur_title, cur_paras))
                cur_title = text
                cur_paras = []
            else:
                cur_paras.append(text)
        if cur_title or cur_paras:
            slides.append(_docx_section_to_slide(cur_title, cur_paras))
    else:
        paras = [t for k, t in blocks if t.strip()]
        if not paras:
            slides.append(
                SampleSlideInfo(
                    title="Документ",
                    bullets=[],
                    body="",
                    word_count=0,
                    has_image=False,
                )
            )
        else:
            for i, chunk in enumerate(_chunk_plain_paragraphs(paras)):
                lines = [ln.strip() for ln in chunk.split("\n") if ln.strip()]
                title_g = lines[0][:120] if lines and len(lines[0]) <= 100 else f"Фрагмент {i + 1}"
                rest = lines[1:] if len(lines) > 1 and len(lines[0]) <= 100 else lines
                body = " ".join(rest).strip()[:1200]
                bullets = [ln[:200] for ln in rest[:10] if len(ln.split()) <= 35]
                info = SampleSlideInfo(
                    title=title_g,
                    bullets=bullets,
                    body=body,
                    word_count=len(chunk.split()),
                    has_image=False,
                )
                slides.append(info)

    n = len(slides) or 1
    for i, info in enumerate(slides):
        info.kind_guess = _guess_kind(i, n, info)

    avg_words = sum(s.word_count for s in slides) / max(1, n)
    density = _density_from_words(avg_words)
    title_guess = slides[0].title if slides else ""

    return SampleAnalysis(
        sample_id=sample_id,
        source_format="docx",
        file_name=path.name,
        n_slides=n,
        palette=dict(_DEFAULT_PALETTE),
        density=density,
        has_images=False,
        slides=slides,
        title_guess=title_guess,
    )


# ============================ entry ============================
def analyze_file(path: Path, sample_id: str) -> SampleAnalysis:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return analyze_pptx(path, sample_id)
    if suffix == ".pdf":
        return analyze_pdf(path, sample_id)
    if suffix == ".docx":
        return analyze_docx(path, sample_id)
    raise ValueError(f"Неподдерживаемый формат: {suffix}. Используйте .pptx, .pdf или .docx")


def sample_outline_for_llm(sample: SampleAnalysis, max_chars: int = 4000) -> str:
    """Готовим компактный JSON-блок для подсказки LLM."""
    is_doc = sample.source_format == "docx"
    if is_doc:
        max_chars = max(max_chars, 12000)
    body_limit = 520 if is_doc else 240
    max_bullets = 10 if is_doc else 6
    payload = {
        "source": sample.source_format,
        "n_slides": sample.n_slides,
        "density": sample.density,
        "palette": sample.palette,
        "has_images": sample.has_images,
        "title_guess": sample.title_guess,
        "slides": [
            {
                "kind": s.kind_guess,
                "title": s.title[:90],
                "bullets": [b[:160] for b in s.bullets[:max_bullets]],
                "body": s.body[:body_limit],
            }
            for s in sample.slides
        ],
    }
    text = json.dumps(payload, ensure_ascii=False)
    if len(text) > max_chars:
        compact = {
            "source": sample.source_format,
            "n_slides": sample.n_slides,
            "density": sample.density,
            "palette": sample.palette,
            "has_images": sample.has_images,
            "title_guess": sample.title_guess,
            "slides": [
                {"kind": s.kind_guess, "title": s.title[:80], "body": s.body[:280]}
                for s in sample.slides
            ],
        }
        text = json.dumps(compact, ensure_ascii=False)
    if len(text) > max_chars:
        compact = {
            "source": sample.source_format,
            "n_slides": sample.n_slides,
            "density": sample.density,
            "palette": sample.palette,
            "title_guess": sample.title_guess[:120],
            "slides": [
                {"kind": s.kind_guess, "title": s.title[:70]}
                for s in sample.slides
            ],
        }
        text = json.dumps(compact, ensure_ascii=False)
    return text
